# utils/pptx_converter.py

from pptx import Presentation
from pathlib import Path
import re

def convert_md_to_pptx(md_content: str, output_pptx: str):
    prs = Presentation()

    # Basic slide layout
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    lines = md_content.split('\n')
    current_title = ""
    current_content = []

    for line in lines:
        if line.startswith("# ") or line.startswith("## "):
            if current_title:
                # Save previous slide
                slide = prs.slides.add_slide(content_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.shapes.placeholders[1]
                title_shape.text = current_title
                tf = body_shape.text_frame
                tf.text = '\n'.join(current_content)
                current_content = []

            current_title = line.lstrip('# ').strip()
        elif line.startswith('- ') or line.startswith('  - '):
            item = line.lstrip('- ').strip()
            current_content.append(item)
        elif line == "---":
            if current_title:
                slide = prs.slides.add_slide(content_slide_layout)
                title_shape = slide.shapes.title
                body_shape = slide.shapes.placeholders[1]
                title_shape.text = current_title
                body_shape.text = '\n'.join(current_content)
                current_title = ""
                current_content = []

    # Save last slide
    if current_title and current_content:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = current_title
        tf = slide.shapes.placeholders[1].text_frame
        for point in current_content:
            p = tf.add_paragraph()
            p.text = point

    # Save presentation
    prs.save(output_pptx)
    print(f"📊 Saved PPTX: {output_pptx}")