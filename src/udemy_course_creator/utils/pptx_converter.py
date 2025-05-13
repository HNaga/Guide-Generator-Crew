from pptx import Presentation
from pptx.util import Inches
from pathlib import Path

def convert_md_to_pptx(md_content: str, output_path: str):
    """
    Convert Markdown-formatted slides into a PowerPoint (.pptx) presentation.
    Supports title, concept, code, and summary slides.
    """
    prs = Presentation()

    # Slide layout helper
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[5]  # Title and Content

    lines = md_content.split('\n')
    current_title = ""
    current_body = []

    for line in lines:
        if line.startswith("# ") or line.startswith("## "):
            if current_title:
                _add_slide(prs, current_title, current_body)
                current_body = []

            current_title = line.lstrip('# ').strip()
        elif line.startswith('- ') or line.strip() == "":
            if line.startswith('- '):
                item = line.lstrip('- ').strip()
                current_body.append(item)
        elif line == "---":
            continue
        else:
            # Handle normal text or headers inside content
            if line.strip():
                current_body.append(line.strip())

    # Add last slide
    if current_title and current_body:
        _add_slide(prs, current_title, current_body)

    # Save file
    prs.save(output_path)
    print(f"📊 Saved PPTX: {output_path}")


def _add_slide(prs, title: str, content: list):
    """Helper to create a slide with title and bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_shape = slide.shapes.title
    body_shape = slide.placeholders[1]

    title_shape.text = title

    tf = body_shape.text_frame
    for i, point in enumerate(content):
        if i == 0:
            tf.text = point
        else:
            p = tf.add_paragraph()
            p.text = point


if __name__ == "__main__":
    sample_md = """
# [Slide 1] Title Slide
## Practical CrewAI Course
### Lecture: What is CrewAI?
> Designed for intermediate learners

---

# [Slide 2] Key Concepts
- CrewAI is a framework for building multi-agent systems
- Agents collaborate using Tasks and Crews

---

# [Slide 3] Code Example
```python
from crewai import Agent, Task, Crew
"""