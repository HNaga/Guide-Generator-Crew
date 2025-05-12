from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
import re
import markdown
from bs4 import BeautifulSoup

# --- ReaderAgent ---
class ReaderAgent:
    def run(self, file_path):
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        html = markdown.markdown(content)
        soup = BeautifulSoup(html, 'html.parser')

        slides = []
        current_title = "Untitled"
        current_content = []

        for tag in soup.contents:
            if tag.name == 'h1':
                if current_content:
                    slides.append({'title': current_title, 'content': current_content})
                current_title = tag.get_text()
                current_content = []
            elif tag.name in ['h2', 'h3']:
                current_content.append(f"**{tag.get_text()}**")
            elif tag.name == 'p':
                current_content.append(tag.get_text())
            elif tag.name == 'ul':
                for li in tag.find_all('li'):
                    current_content.append(f"- {li.get_text()}")
            elif tag.name == 'ol':
                for i, li in enumerate(tag.find_all('li'), start=1):
                    current_content.append(f"{i}. {li.get_text()}")
            elif tag.name == 'blockquote':
                current_content.append(f"> {tag.get_text()}")
            elif tag.name == 'pre':
                current_content.append(f"`{tag.get_text()}`")

        if current_content:
            slides.append({'title': current_title, 'content': current_content})
        return slides

# --- SlidePlanner ---
class SlidePlanner:
    def __init__(self, max_chars=600):
        self.max_chars = max_chars

    def run(self, parsed_sections):
        all_slides = []
        for section in parsed_sections:
            title = section['title']
            content = section['content']
            slides = self.split_by_char_limit(content, self.max_chars)
            for i, chunk in enumerate(slides):
                slide_title = f"{title} (Part {i+1})" if len(slides) > 1 else title
                all_slides.append({'title': slide_title, 'content': chunk})
        return all_slides

    def split_by_char_limit(self, lines, max_chars):
        slides = []
        current = []
        count = 0
        for line in lines:
            line_length = len(line)
            if count + line_length > max_chars:
                slides.append(current)
                current = [line]
                count = line_length
            else:
                current.append(line)
                count += line_length
        if current:
            slides.append(current)
        return slides

# --- SlideWriter ---
class SlideWriter:
    def run(self, slides, output_file='presentation.pptx'):
        prs = Presentation()
        slide_layout = prs.slide_layouts[1]  # title + content

        for slide_data in slides:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = slide_data['title']
            tf = slide.placeholders[1].text_frame
            tf.clear()

            for line in slide_data['content']:
                para = tf.add_paragraph()
                para.text = self.clean_md(line)
                para.level = 0
                para.font.size = Pt(20)

        prs.save(output_file)
        print(f"✅ Presentation saved as {output_file}")

    def clean_md(self, line):
        # Strip simple markdown formatting
        line = re.sub(r'\*\*(.*?)\*\*', r'\1', line)
        line = re.sub(r'\*(.*?)\*', r'\1', line)
        line = re.sub(r'`(.*?)`', r'\1', line)
        line = re.sub(r'^- ', '• ', line)
        line = re.sub(r'^> ', '❝ ', line)
        return line

# --- Main Crew Orchestration ---
def main():
    file_path = 'output\CrewAI 101- Introduction to Autonomous AI Agents.md'  # Replace with your file

    reader = ReaderAgent()
    planner = SlidePlanner()
    writer = SlideWriter()

    sections = reader.run(file_path)
    slides = planner.run(sections)
    writer.run(slides, 'output\CrewAI 101- Introduction to Autonomous AI Agents.pptx')

if __name__ == "__main__":
    main()
